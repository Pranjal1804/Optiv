import os
import re
import shutil
import tempfile
import json
from dotenv import load_dotenv
import csv
import base64
from typing import Union, Dict
from datetime import datetime
import pdfplumber
import docx
from pptx import Presentation
import openpyxl
import xlrd
from PIL import Image
import pytesseract
import torch
from transformers import pipeline
from langchain.agents import initialize_agent, Tool
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.messages import HumanMessage
from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
import uvicorn
import io
import subprocess
import zipfile

load_dotenv()
pytesseract.pytesseract.tesseract_cmd = "/usr/bin/tesseract"
os.environ['TESSDATA_PREFIX'] = '/usr/share/tessdata'

def detect_file_type(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf": return "pdf"
    if ext in [".docx", ".doc"]: return "docx"
    if ext == ".pptx": return "pptx"
    if ext == ".xlsx": return "xlsx"
    if ext == ".xls": return "xls"
    if ext == ".csv": return "csv"
    if ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]: return "image"
    return "unknown"

def clean_text(text):
    if not text: return ""
    text = "\n".join([line.strip() for line in text.splitlines() if line.strip()])
    text = re.sub(r"[^\x20-\x7E\n]", " ", text)
    return text.strip()

def extract_pdf_text(file_path):
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text(x_tolerance=2, y_tolerance=2)
                if page_text and page_text.strip():
                    text += page_text + "\n"
                else:
                    try:
                        page_image = page.to_image(resolution=300).original
                        ocr_text = pytesseract.image_to_string(page_image)
                        text += ocr_text + "\n"
                    except Exception:
                        pass
    except Exception as e:
        print(f"PDF extraction error: {e}")
    return clean_text(text)

def extract_docx_text(file_path):
    text = ""
    try:
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip(): text += para.text + "\n"
    except Exception as e:
        print(f"DOCX extraction error: {e}")
    return clean_text(text)

def extract_pptx_text(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip(): text += shape.text + "\n"
    except Exception as e:
        print(f"PPTX extraction error: {e}")
    return clean_text(text)

def extract_xlsx_text(file_path):
    text = ""
    try:
        if not zipfile.is_zipfile(file_path):
            return "Error reading Excel file: not a valid XLSX (zip) file"
    except Exception:
        pass
    try:
        tmpdir = tempfile.mkdtemp(prefix="xlsx2csv_")
        csv_paths = []
        if shutil.which("ssconvert"):
            cmd = ["ssconvert", "--export-type=Gnumeric_stf:stf_csv", "--export-file-per-sheet", file_path, os.path.join(tmpdir, "out.csv")]
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            for f in sorted(os.listdir(tmpdir)):
                if f.endswith(".csv"):
                    csv_paths.append(os.path.join(tmpdir, f))
        elif shutil.which("soffice"):
            cmd = ["soffice", "--headless", "--convert-to", "csv", "--outdir", tmpdir, file_path]
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            for f in sorted(os.listdir(tmpdir)):
                if f.endswith(".csv"):
                    csv_paths.append(os.path.join(tmpdir, f))
        else:
            return "Error reading Excel file: no converter (ssconvert/soffice) available"
        collected = []
        for p in csv_paths:
            collected.append(f"--- Sheet CSV: {os.path.basename(p)} ---")
            collected.append(extract_csv_text(p))
        text = "\n".join(collected)
        return clean_text(text) if text.strip() else "No readable content found in Excel file"
    except Exception as e:
        print(f"XLSX CLI conversion error: {e}")
    try:
        import pandas as pd
        try:
            sheets = pd.read_excel(file_path, sheet_name=None, engine="openpyxl")
        except Exception:
            sheets = pd.read_excel(file_path, sheet_name=None)
        out = []
        for name, df in sheets.items():
            out.append(f"--- Sheet: {name} ---")
            if df is not None and not df.empty:
                df = df.fillna("")
                for _, row in df.iterrows():
                    row_vals = [str(x).strip() for x in row.tolist() if str(x).strip() != ""]
                    if row_vals:
                        out.append(" | ".join(row_vals))
        text = "\n".join(out)
        return clean_text(text) if text.strip() else "No readable content found in Excel file"
    except Exception as e:
        print(f"XLSX pandas/calamine error: {e}")
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True, keep_links=False)
        for sheet in wb.worksheets:
            text += f"--- Sheet: {sheet.title} ---\n"
            for row in sheet.iter_rows(values_only=True):
                if not row:
                    continue
                vals = []
                for v in row:
                    if v is None:
                        continue
                    if isinstance(v, datetime):
                        vals.append(v.strftime("%Y-%m-%d %H:%M:%S"))
                    else:
                        vals.append(str(v).strip())
                if vals:
                    line = " | ".join([v for v in vals if v])
                    if line.strip():
                        text += line + "\n"
        wb.close()
        return clean_text(text) if text.strip() else "No readable content found in Excel file"
    except Exception as e:
        print(f"XLSX openpyxl error: {e}")
        return f"Error reading Excel file: {str(e)}"

def extract_xls_text(file_path):
    text = ""
    try:
        wb = xlrd.open_workbook(file_path)
        for sheet_idx in range(wb.nsheets):
            sheet = wb.sheet_by_index(sheet_idx)
            text += f"--- Sheet: {sheet.name} ---\n"
            for row_idx in range(sheet.nrows):
                row_values = []
                for col_idx in range(sheet.ncols):
                    try:
                        cell = sheet.cell(row_idx, col_idx)
                        cell_value = cell.value
                        if cell.ctype == xlrd.XL_CELL_DATE:
                            dt_tuple = xlrd.xldate_as_tuple(cell_value, wb.datemode)
                            cell_value = datetime(*dt_tuple).strftime("%Y-%m-%d %H:%M:%S")
                        if cell_value is not None and str(cell_value).strip():
                            row_values.append(str(cell_value))
                    except Exception:
                        continue
                if row_values:
                    row_text = " | ".join(row_values)
                    if row_text.strip():
                        text += row_text + "\n"
    except Exception as e:
        print(f"XLS extraction error: {e}")
        return f"Error reading Excel file: {str(e)}"
    return clean_text(text) if text.strip() else "No readable content found in Excel file"

def extract_csv_text(file_path):
    text = ""
    try:
        with open(file_path, mode='r', encoding='utf-8', errors='ignore') as csvfile:
            reader = csv.reader(csvfile)
            for row in reader:
                row_text = " | ".join([cell for cell in row if cell and cell.strip() != ""])
                if row_text.strip():
                    text += row_text + "\n"
    except Exception as e:
        print(f"CSV extraction error: {e}")
        return f"Error reading CSV file: {str(e)}"
    return clean_text(text)

def extract_image_text(file_path):
    text = ""
    try:
        img = Image.open(file_path)
        img = img.convert("L")
        text = pytesseract.image_to_string(img)
    except Exception as e:
        print(f"Image OCR error: {e}")
    return clean_text(text)

def extract_text_from_file(file_path):
    ftype = detect_file_type(file_path)
    if ftype == "pdf": return extract_pdf_text(file_path), "PDF"
    elif ftype == "docx": return extract_docx_text(file_path), "DOCX"
    elif ftype == "pptx": return extract_pptx_text(file_path), "PPTX"
    elif ftype == "xlsx": return extract_xlsx_text(file_path), "XLSX"
    elif ftype == "xls": return extract_xls_text(file_path), "XLS"
    elif ftype == "csv": return extract_csv_text(file_path), "CSV"
    elif ftype == "image": return extract_image_text(file_path), "Image"
    else:
        try:
            with open(file_path, "r", errors="ignore") as f: return clean_text(f.read()), "Text"
        except:
            return "", "Unknown"

def convert_excel_to_pdf_or_docx(file_path: str):
    try:
        if shutil.which("soffice"):
            outdir = tempfile.mkdtemp(prefix="excel2pdf_")
            subprocess.run(
                ["soffice", "--headless", "--convert-to", "pdf", "--outdir", outdir, file_path],
                check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE
            )
            base = os.path.splitext(os.path.basename(file_path))[0] + ".pdf"
            pdf_path = os.path.join(outdir, base)
            if os.path.exists(pdf_path):
                return pdf_path, "pdf"
    except Exception as e:
        print(f"Excel->PDF (soffice) error: {e}")
    return None, None

ner = pipeline("ner", model="dslim/bert-base-NER", grouped_entities=True, device=0 if torch.cuda.is_available() else -1)
summarizer = pipeline("summarization", model="facebook/bart-large-cnn", device=0 if torch.cuda.is_available() else -1)

REGEX_PATTERNS = {
    "EMAIL": re.compile(r"\b[a-zA-Z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b"),
    "IP_ADDRESS": re.compile(r"\b(?:\d{1,3}\.){3}\d{1,3}\b"),
    "PHONE": re.compile(r"\b(?:\+?\d{1,3}[-.\s]?)?(?:\(?\d{2,4}\)?[-.\s]?)?\d{3,4}[-.\s]?\d{3,4}\b"),
    "API_KEY": re.compile(r"\b(api_key|apikey|api-key|token|secret|password)\s*[:=]\s*['\"]?([a-zA-Z0-9\-_]{20,})['\"]?\b", re.IGNORECASE),
    "PRIVATE_KEY": re.compile(r"-----BEGIN (RSA|EC|OPENSSH|PGP) PRIVATE KEY-----"),
}

def mask_structured_pii(text):
    found = set()
    out = text
    for label, pat in REGEX_PATTERNS.items():
        for m in pat.finditer(out):
            snippet = m.group(0)
            found.add(label)
            out = out.replace(snippet, f"[REDACTED_{label}]")
    return out, list(found)

def mask_ner_pii(text):
    if not text or not text.strip(): return text, []
    try:
        ents = ner(text)
    except Exception:
        ents = []
    if not ents: return text, []
    out_chars = list(text)
    for ent in sorted(ents, key=lambda x: x['start'], reverse=True):
        out_chars[ent['start']:ent['end']] = list(f"[REDACTED_{ent['entity_group']}]")
    out_text = "".join(out_chars)
    labels_found = list({ent['entity_group'] for ent in ents})
    return out_text, labels_found

def extract_security_insights(text):
    findings = {}
    fw_pattern = re.compile(r"\b(ALLOW|DENY)\b.*?(?:\b(?:\d{1,3}\.){3}\d{1,3}\b(?::\d+)?|\bport\s+\d+\b|\bport:\s*\d+\b)", flags=re.IGNORECASE)
    for m in fw_pattern.finditer(text):
        findings.setdefault("FIREWALL_RULE", []).append(m.group(0).strip())
    iam_pattern = re.compile(r'("Effect"\s*:\s*"(Allow|Deny)"|IAM\s+policy|"Action"\s*:)', flags=re.IGNORECASE)
    for m in iam_pattern.finditer(text):
        start = max(0, m.start()-50)
        end = min(len(text), m.end()+150)
        findings.setdefault("IAM_POLICY_MENTION", []).append(text[start:end].strip())
    port_pattern = re.compile(r"\bport[:\s]*\d{1,5}\b", flags=re.IGNORECASE)
    for m in port_pattern.finditer(text):
        findings.setdefault("PORT", []).append(m.group(0))
    for label, pat in REGEX_PATTERNS.items():
        matches = sorted(set(pat.findall(text)))
        if matches: findings.setdefault(f"Potential {label}s", []).extend(matches)
    vuln_keywords = re.compile(r"\b(sql injection|xss|cross-site scripting|csrf|rce|remote code execution|password|secret|vulnerability|exploit)\b", re.IGNORECASE)
    vuln_mentions = sorted(set(vuln_keywords.findall(text)))
    if vuln_mentions: findings["Security Keyword Mentions"] = vuln_mentions
    return findings

def summarize_text(text, chunk_size=800, max_length=150, min_length=40):
    if not text.strip(): return "No content to summarize"
    words = text.split()
    chunks = [" ".join(words[i:i+chunk_size]) for i in range(0, len(words), chunk_size)]
    summaries = [summarizer(chunk, max_length=max_length, min_length=min_length, do_sample=False)[0]['summary_text'] for chunk in chunks]
    return " ".join(summaries)

def generate_descriptive_output(file_name, file_type, summary, insights, pii_removed):
    key_findings = [f"{k}: {', '.join(map(str, v[:3]))}" for k, v in insights.items()]
    return {
        "File Name": file_name,
        "File Type": file_type,
        "File Description": summary,
        "PII Removed": ", ".join(pii_removed) if pii_removed else "None",
        "Key Findings": "; ".join(key_findings) if key_findings else "No major security insights detected"
    }

api_key = os.environ.get("GOOGLE_API_KEY")
if not api_key: raise ValueError("GOOGLE_API_KEY environment variable not set.")
llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash", google_api_key=api_key)

def image_to_base64(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

def describe_image_content(image_path: str) -> str:
    try:
        b64_string = image_to_base64(image_path)
        message = HumanMessage(
            content=[
                {"type": "text", "text": "Describe this image in detail. What objects are present? What is the setting? What is happening?"},
                {"type": "image_url","image_url": f"data:image/png;base64,{b64_string}"},
            ]
        )
        response = llm.invoke([message])
        return response.content
    except Exception as e:
        return f"Could not describe image due to an error: {e}"

def descriptive_tool(tool_input: Union[str, Dict[str, str]]):
    if isinstance(tool_input, dict):
        path = tool_input.get("path") or tool_input.get("input")
    else:
        path = tool_input
    if not path:
        return {"error": "No file path provided to the tool."}
    raw_text, ftype = extract_text_from_file(path)
    if ftype == "Image" and (not raw_text or len(raw_text) < 25):
        description = describe_image_content(path)
        return {
            "File Name": os.path.basename(path),
            "File Type": "Image (Visual Analysis)",
            "File Description": description,
            "PII Removed": "N/A",
            "Key Findings": "N/A (Visual analysis performed)"
        }
    if not raw_text: return {"error": "Could not extract text from file."}
    insights = extract_security_insights(raw_text)
    ner_masked_text, ner_labels = mask_ner_pii(raw_text)
    final_text, structured_labels = mask_structured_pii(ner_masked_text)
    pii_removed = sorted(list(set(ner_labels + structured_labels)))
    summary = summarize_text(final_text)
    record = generate_descriptive_output(os.path.basename(path), ftype, summary, insights, pii_removed)
    return record

tools = [
    Tool(
        name="DescriptiveOutput",
        func=descriptive_tool,
        description="Analyzes a file at a given path, masks PII, and generates a structured output in a dictionary format. Use this tool for any file analysis request."
    ),
]

agent = initialize_agent(
    tools,
    llm,
    agent="zero-shot-react-description",
    verbose=True,
    handle_parsing_errors=True
)

app = FastAPI(title="OptivSec File Analysis API", version="1.4")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

@app.post("/analyze_file/")
async def analyze_file_endpoint(file: UploadFile = File(...)):
    temp_file_path = ""
    conv_path = None
    try:
        file_bytes = await file.read()  
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.filename)[1]) as tmp_file:
            tmp_file.write(file_bytes)
            tmp_file.flush() 
            os.fsync(tmp_file.fileno())  
            temp_file_path = tmp_file.name

        if file.filename.lower().endswith(('.xlsx', '.xls')):
            try:
                if file.filename.lower().endswith('.xlsx'):
                    if not zipfile.is_zipfile(temp_file_path):
                        return {
                            "File Name": file.filename,
                            "File Type": "XLSX",
                            "File Description": "Not a valid XLSX file (not a zip archive)",
                            "PII Removed": "None",
                            "Key Findings": "Invalid file format"
                        }
                elif file.filename.lower().endswith('.xls'):
                    try:
                        xlrd.open_workbook(temp_file_path)
                    except Exception:
                        return {
                            "File Name": file.filename,
                            "File Type": "XLS",
                            "File Description": "Not a valid XLS file",
                            "PII Removed": "None",
                            "Key Findings": "Invalid file format"
                        }

                conv_path, conv_kind = convert_excel_to_pdf_or_docx(temp_file_path)
                if not conv_path:
                    raw_text, ftype = extract_text_from_file(temp_file_path)
                    if isinstance(raw_text, str) and raw_text.startswith("Error reading Excel file:"):
                        return {
                            "File Name": file.filename,
                            "File Type": "XLSX",
                            "File Description": raw_text,
                            "PII Removed": "None",
                            "Key Findings": "Processing failed"
                        }
                    if not raw_text or len(raw_text.strip()) < 5:
                        return {
                            "File Name": file.filename,
                            "File Type": ftype,
                            "File Description": "Empty or unreadable Excel file",
                            "PII Removed": "None",
                            "Key Findings": "No content found"
                        }
                    insights = extract_security_insights(raw_text)
                    ner_masked_text, ner_labels = mask_ner_pii(raw_text)
                    final_text, structured_labels = mask_structured_pii(ner_masked_text)
                    pii_removed = sorted(list(set(ner_labels + structured_labels)))
                    summary = summarize_text(final_text)
                    return generate_descriptive_output(
                        file.filename,
                        ftype,
                        summary,
                        insights,
                        pii_removed
                    )

                if conv_kind == "pdf":
                    raw_text = extract_pdf_text(conv_path)
                    ftype = "PDF (from Excel)"
                else:
                    raw_text = extract_docx_text(conv_path)
                    ftype = "DOCX (from Excel)"

                if not raw_text or len(raw_text.strip()) < 5:
                    return {
                        "File Name": file.filename,
                        "File Type": ftype,
                        "File Description": "Empty or unreadable content after conversion",
                        "PII Removed": "None",
                        "Key Findings": "No content found"
                    }

                insights = extract_security_insights(raw_text)
                ner_masked_text, ner_labels = mask_ner_pii(raw_text)
                final_text, structured_labels = mask_structured_pii(ner_masked_text)
                pii_removed = sorted(list(set(ner_labels + structured_labels)))
                summary = summarize_text(final_text)

                return generate_descriptive_output(
                    file.filename,
                    ftype,
                    summary,
                    insights,
                    pii_removed
                )
            except Exception as excel_error:
                return {
                    "File Name": file.filename,
                    "File Type": "XLSX",
                    "File Description": f"Error processing Excel: {str(excel_error)}",
                    "PII Removed": "None",
                    "Key Findings": "Processing failed"
                }
            finally:
                if conv_path and os.path.exists(conv_path):
                    try:
                        os.remove(conv_path)
                    except Exception:
                        pass

        prompt = f"Analyze the file at the following path and provide the structured output: {temp_file_path}"
        response = agent.run(prompt)
        try:
            dict_match = re.search(r"\{.*\}", response, re.DOTALL)
            if dict_match:
                dict_str = dict_match.group()
                import ast
                analysis_result = ast.literal_eval(dict_str)
            else:
                raise ValueError(f"Agent did not return a structured result: {response}")
            return analysis_result
        except (SyntaxError, ValueError) as e:
            raise HTTPException(status_code=500, detail=f"Failed to parse agent response: {response}. Error: {e}")

    except Exception as e:
        if isinstance(e, HTTPException): raise e
        raise HTTPException(status_code=500, detail=f"An error occurred during analysis: {str(e)}")
    finally:
        if temp_file_path and os.path.exists(temp_file_path):
            os.remove(temp_file_path)

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)